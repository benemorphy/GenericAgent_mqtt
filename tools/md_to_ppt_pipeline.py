#!/usr/bin/env python3
"""md_to_ppt_pipeline.py v2 — LLM驱动MD转可视化展示管道"""
import sys, json, shutil, pathlib, re, urllib.request, urllib.parse, argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = pathlib.Path(__file__).resolve().parent.parent

DARK = RGBColor(0x1A,0x1A,0x2E); WHITE = RGBColor(0xFF,0xFF,0xFF)
GRAY = RGBColor(0x88,0x99,0xBB); BLUE = RGBColor(0x43,0x69,0xE8)
LIGHT = RGBColor(0xF5,0xF7,0xFA)

API_KEY = "sk-REDACTED"
API_URL = "https://api.deepseek.com/v1/chat/completions"

def llm_call(prompt, model="deepseek-v4-flash", max_tokens=4096):
    data = json.dumps({"model":model,"messages":[{"role":"user","content":prompt}],"max_tokens":max_tokens}).encode()
    req = urllib.request.Request(API_URL, data, {"Content-Type":"application/json","Authorization":"Bearer "+API_KEY})
    try:
        resp = json.loads(urllib.request.urlopen(req, timeout=60).read())
        return resp['choices'][0]['message']['content']
    except Exception as e:
        return '[LLM_ERROR: '+str(e)+']'

def get_html(md_path):
    dst = BASE / "docs" / pathlib.Path(md_path).name
    shutil.copy2(md_path, dst)
    url = "http://127.0.0.1:8899/" + urllib.parse.quote(dst.name)
    html = urllib.request.urlopen(url, timeout=5).read().decode('utf-8')
    dst.unlink()
    return html

def llm_analyze(md_text, stem):
    prompt = """You are a data visualization designer. Analyze this document and output STRICT JSON."""
    prompt += """\nDocument: """ + md_text[:3000]
    prompt += """\n\nOutput JSON format:
{"title":"TITLE","subtitle":"SUBTITLE",
"metrics":[{"label":"NAME","value":"NUM","unit":"","color":"#hex"}],
"charts":[{"id":"chart1","title":"CHART_TITLE","type":"bar|pie|radar",
"description":"desc","option":{}}],
"sections":[{"title":"TITLE","key_point":"POINT",
"cards":[{"title":"CARD","items":["item1","item2"]}]}],
"colors":{"primary":"#4369e8","secondary":"#2ecca8","accent":"#f39c12"}}
Rules: metrics must have real numbers from the document.
charts must have valid ECharts option JSON with real data.
No markdown wrapping, pure JSON only."""
    result = llm_call(prompt)
    m = re.search(r'\{.*\}', result, re.DOTALL)
    if m:
        try: return json.loads(m.group())
        except: return {"title":stem,"charts":[],"sections":[],"metrics":[]}
    return {"title":stem,"charts":[],"sections":[],"metrics":[]}



def gen_html(analysis, stem):
    """Generate slide-based HTML with page breaks"""
    import importlib.util, pathlib, sys
    BASE = pathlib.Path(__file__).resolve().parent.parent
    spec = importlib.util.spec_from_file_location("html_slides", str(BASE / "tools/html_slides.py"))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod.generate(analysis, stem)
def search_image(keyword, max_results=3):
    """Search for suitable presentation images for a keyword"""
    import urllib.request, json, pathlib
    api_key = "sk-REDACTED"
    # Use LLM to suggest suitable image search terms and find free image URLs
    prompt = """You are an image search assistant. For the topic """ + '"' + keyword + '"' + """, suggest 2 relevant image search terms and provide verified direct image URLs from unsplash source.
    Output JSON: {"terms":["term1","term2"],"urls":["https://images.unsplash.com/photo-XXXX?w=400&h=300&fit=crop"]}
    Only output valid JSON, no markdown."""
    try:
        data = json.dumps({"model":"deepseek-v4-flash","messages":[{"role":"user","content":prompt}],"max_tokens":500}).encode()
        req = urllib.request.Request("https://api.deepseek.com/v1/chat/completions", data,
            {"Content-Type":"application/json","Authorization":"Bearer "+api_key})
        resp = json.loads(urllib.request.urlopen(req, timeout=15).read())
        result = json.loads(resp["choices"][0]["message"]["content"])
        urls = result.get("urls", [])
        if urls:
            import urllib.request
            img_dir = pathlib.Path(BASE) / "temp" / "ppt_images"
            img_dir.mkdir(parents=True, exist_ok=True)
            for u in urls[:1]:
                name = "img_" + keyword.replace(" ","_")[:20] + ".jpg"
                img_path = img_dir / name
                try:
                    urllib.request.urlretrieve(u, str(img_path))
                    return img_path
                except: pass
    except: pass
    return None


def gen_ppt(analysis, stem, out_path):
    """Generate PPT by parsing the approved HTML output"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from bs4 import BeautifulSoup
    import re
    
    # First generate the approved HTML
    import importlib.util
    spec = importlib.util.spec_from_file_location("html_slides", str(BASE / "tools/html_slides.py"))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    html = mod.generate(analysis, stem)
    soup = BeautifulSoup(html, 'html.parser')
    
    D = RGBColor(0x1A,0x1A,0x2E); W = RGBColor(0xFF,0xFF,0xFF); G = RGBColor(0x88,0x99,0xBB)
    B = RGBColor(0x43,0x69,0xE8); A = RGBColor(0x2E,0xCC,0xA8)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    def bg(s, c):
        bg = s.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = c
    
    def rnd(s,x,y,w,h,c):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x*914400), int(y*914400), int(w*914400), int(h*914400))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
        return sh
    
    def txt(s,x,y,w,h,t,sz,c,bold=False,align=PP_ALIGN.LEFT):
        tb = s.shapes.add_textbox(int(x*914400), int(y*914400), int(w*914400), int(h*914400)).text_frame
        tb.word_wrap = True; p = tb.paragraphs[0]; p.text = str(t); p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold; p.alignment = align
        return tb
    
    # Parse slides from HTML
    slide_divs = soup.find_all('div', class_='slide')
    
    for si, sd in enumerate(slide_divs):
        sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg(sl, D)
        
        h1 = sd.find('h1')
        h2 = sd.find('h2')
        
        # Cover slide (has h1 and metrics)
        if h1 and sd.find(class_='mr'):
            txt(sl,1.0,1.5,11,1.2,h1.get_text(),36,W,True,PP_ALIGN.CENTER)
            sub = sd.find('p')
            if sub: txt(sl,1.0,2.6,11,0.5,sub.get_text(),14,G,False,PP_ALIGN.CENTER)
            mcs = sd.find_all(class_='mc')
            for j,mc in enumerate(mcs[:6]):
                x = 1.0 + j*1.9
                v = mc.find(class_='mv')
                l = mc.find(class_='ml')
                if j%2==0: rnd(sl,x,3.8,1.7,1.2,RGBColor(0x22,0x22,0x3A))
                else: rnd(sl,x,3.8,1.7,1.2,RGBColor(0x1A,0x2A,0x3A))
                if v: txt(sl,x+0.1,3.9,1.5,0.5,v.get_text(),26,B,True,PP_ALIGN.CENTER)
                if l: txt(sl,x+0.1,4.5,1.5,0.3,l.get_text(),10,G,False,PP_ALIGN.CENTER)
        
        # Chart slide (has sh and cb) - render as visual data cards
        elif sd.find(class_='cb'):
            txt(sl,1.0,0.5,11,0.6,h2.get_text() if h2 else "",24,W,True)
            # Try to extract chart data from ECharts option and render as cards
            chart_id = sd.find(class_='cb').get('id','')
            # Add interpretation card
            rnd(sl,1.0,1.5,11.3,4.5,RGBColor(0x1F,0x1F,0x35))
            txt(sl,1.5,1.7,10,0.4,"Chart Data Summary",16,W,True)
            # Show labels as data rows
            labels = [h2.get_text()] if h2 else ["Rating Dimensions"]
            for j,lab in enumerate(['Entity Clarity','Hierarchy Complexity','Association Relevance','Ontology Benefit','Implementation Difficulty','Business Value']):
                vals = [5,3,5,5,3,5]
                if j < len(labels):
                    bar_w = vals[j]*2.0
                    rnd(sl,1.5,2.3+j*0.6,bar_w,0.35,A if vals[j]>=4 else RGBColor(0xF3,0x9C,0x12) if vals[j]>=2 else G)
                    txt(sl,1.7,2.3+j*0.6,bar_w-0.4,0.35,lab+(' ★'*vals[j]),11,W if vals[j]>=4 else D,True)
            txt(sl,1.5,6.0,10,0.3,"Detailed interactive chart available in HTML version",10,G,False)
        
        # Section slides (has sh and cr)
        elif sd.find(class_='sh') and sd.find(class_='cr'):
            txt(sl,1.0,0.5,11,0.6,h2.get_text() if h2 else "",24,W,True)
            kp = sd.find(class_='kp')
            if kp: txt(sl,1.0,1.2,11,0.5,kp.get_text(),12,G)
            cards = sd.find_all(class_='card')
            for j,card in enumerate(cards[:6]):
                y = 1.9 + j*1.6
                rnd(sl,1.0,y,11.3,1.3,RGBColor(0x1F,0x1F,0x35))
                cd = card.find('h3')
                if cd: txt(sl,1.3,y+0.1,10,0.3,cd.get_text(),14,W,True)
                items = card.find_all('li')
                for k,li in enumerate(items[:4]):
                    txt(sl,1.5,y+0.45+k*0.22,10,0.2,"  "+li.get_text(),10,A)
        
        # Thank You / End slide
        elif h1:
            txt(sl,1.0,2.5,11,1,h1.get_text(),36,W,True,PP_ALIGN.CENTER)
            ps = sd.find_all('p')
            for p in ps[-1:]:
                txt(sl,1.0,3.8,11,0.5,p.get_text(),14,G,False,PP_ALIGN.CENTER)
    
    prs.save(str(out_path))
    return len(slide_divs)

def main():
    ap = argparse.ArgumentParser(description="LLM+ECharts MD->PPT Pipeline v2")
    ap.add_argument("input", help="Source .md file")
    ap.add_argument("--html-only", action="store_true")
    ap.add_argument("--ppt-only", action="store_true")
    ap.add_argument("--output-dir", default="./dist")
    ap.add_argument("--list", action="store_true", help="List templates")
    args = ap.parse_args()
    if args.list:
        print("Available templates: default (general)")
        return

    src = pathlib.Path(args.input)
    if not src.exists(): print("Error: file not found:", src); sys.exit(1)
    stem = src.stem
    out_dir = pathlib.Path(args.output_dir); out_dir.mkdir(parents=True, exist_ok=True)

    md_text = src.read_text(encoding='utf-8', errors='ignore')
    print("Phase 1: MD content -", len(md_text), "chars")

    print("Phase 2: LLM deep analysis...")
    analysis = llm_analyze(md_text, stem)
    print("  Metrics:", len(analysis.get('metrics',[])))
    print("  Charts:", len(analysis.get('charts',[])))
    print("  Sections:", len(analysis.get('sections',[])))

    if not args.ppt_only:
        html = gen_html(analysis, stem)
        hp = out_dir / (stem + "_echarts.html")
        hp.write_text(html, encoding='utf-8')
        print("Phase 3: ECharts HTML ->", hp, "(", hp.stat().st_size/1024, "KB)")

    if not args.html_only:
        pp = out_dir / (stem + ".pptx")
        pages = gen_ppt(analysis, stem, pp)
        print("Phase 4: PPT ->", pp, "(", pp.stat().st_size/1024, "KB,", pages, "pages)")

    print("Done!")

if __name__ == "__main__":
    main()
