#!/usr/bin/env python3
"""
echart_ppt_pipeline.py — ECharts HTML预览 → pyecharts转换 → PPT生成管道
========================================================================
三阶段+HITL(人工介入)工作流:

  preview  : ECharts option JSON → 独立HTML (浏览器预览)
  convert  : option JSON → pyecharts Python代码 (LLM辅助转换)
  build    : pyecharts代码 → PNG截图 → python-pptx组装PPT
  run      : 完整管道 (preview → Gate1 → convert → Gate2 → build)

用法:
  python tools/echart_ppt_pipeline.py preview <chart.json> [--out-dir ./dist]
  python tools/echart_ppt_pipeline.py convert <chart.json> [--out-dir ./dist]
  python tools/echart_ppt_pipeline.py build <presentation.json> [--out-dir ./dist]
  python tools/echart_ppt_pipeline.py run <presentation.json> [--out-dir ./dist]

JSON格式示例:
  // chart.json — 单图表
  {
    "slide_title": "Revenue",
    "chart_title": "Monthly Revenue",
    "width": 800,
    "height": 450,
    "option": {
      "xAxis": {"type": "category", "data": ["Jan","Feb","Mar"]},
      "yAxis": {"type": "value"},
      "series": [{"type": "line", "data": [100, 200, 150]}]
    }
  }
  
  // presentation.json — 完整演示文稿
  {
    "title": "Q1 Report",
    "theme": {"primary": "#4369e8", "secondary": "#2ecca8", "accent": "#f39c12"},
    "slides": [
      {"type": "cover", "title": "Q1 2026 Report", "subtitle": "Data Overview",
       "metrics": [{"label":"Revenue","value":"$1.2M"}, {"label":"Users","value":"50K"}]},
      {"type": "chart", "title": "Revenue Trend", "option": {...}},
      {"type": "content", "title": "Key Insights", "key_point": "...", 
       "cards": [{"title":"Card","items":["item1","item2"]}]},
      {"type": "end", "title": "Thank You"}
    ]
  }
"""

import sys, json, pathlib, argparse, time, os, re, textwrap, shutil

BASE = pathlib.Path(__file__).resolve().parent.parent
DIST_DIR = BASE / "dist"

# ─────────────────────────────────────────────
# Phase 1: ECharts option JSON → HTML 预览
# ─────────────────────────────────────────────

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<title>{title}</title>
<script src="https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js"></script>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#1a1a2e;display:flex;justify-content:center;align-items:center;min-height:100vh;font-family:'Microsoft YaHei',sans-serif}}
.chart-container{{width:{width}px;height:{height}px;background:#fff;border-radius:12px;padding:20px;box-shadow:0 8px 32px rgba(0,0,0,0.3)}}
.chart-title{{font-size:18px;font-weight:600;color:#333;margin-bottom:12px;text-align:center}}
#chart{{width:100%;height:{inner_height}px}}
.controls{{position:fixed;bottom:20px;left:50%;transform:translateX(-50%);display:flex;gap:12px}}
.controls a{{padding:8px 20px;border-radius:20px;background:#4369e8;color:#fff;text-decoration:none;font-size:13px}}
</style>
</head>
<body>
<div class="chart-container">
<div class="chart-title">{chart_title}</div>
<div id="chart"></div>
</div>
<div class="controls">
<a href="#" onclick="var c=document.getElementById('chart');c.style.width=parseInt(c.style.width||'{width}')*1.2+'px';c.resize()">放大</a>
<a href="#" onclick="var c=document.getElementById('chart');c.style.width='100%';c.resize()">重置</a>
</div>
<script>
var chart = echarts.init(document.getElementById('chart'));
chart.setOption({option_json});
window.addEventListener('resize', function() {{ chart.resize() }});
</script>
</body>
</html>"""


def cmd_preview(args):
    """Phase 1: 从ECharts option JSON生成预览HTML"""
    data = _load_json(args.input)
    out_dir = pathlib.Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    
    # 支持单chart或presentation格式
    if "slides" in data:
        # presentation格式：为每个chart slide生成HTML
        outputs = []
        for i, slide in enumerate(data.get("slides", [])):
            if slide.get("type") == "chart" and "option" in slide:
                html = _render_chart_html(slide, data.get("theme", {}))
                stem = slide.get("chart_title", slide.get("title", f"chart_{i}"))
                safe = _safe_name(stem)
                out_path = out_dir / f"{safe}_preview.html"
                out_path.write_text(html, encoding="utf-8")
                outputs.append(out_path)
        if outputs:
            print(f"[OK] 生成 {len(outputs)} 个预览HTML:")
            for p in outputs:
                print(f"   file:///{p.as_posix()}")
        else:
            print("[WARN]  未找到任何 chart slide")
    else:
        # 单chart格式
        html = _render_chart_html(data)
        stem = data.get("chart_title", data.get("slide_title", "chart"))
        out_path = out_dir / f"{_safe_name(stem)}_preview.html"
        out_path.write_text(html, encoding="utf-8")
        print(f"[OK] 预览HTML: file:///{out_path.as_posix()}")
    
    return outputs if "slides" in data else [out_path]


def _render_chart_html(data, theme=None):
    """渲染单个ECharts图表为HTML"""
    w = data.get("width", 800)
    h = data.get("height", 450)
    title = data.get("chart_title", data.get("slide_title", "Chart"))
    option = data.get("option", {})
    
    return HTML_TEMPLATE.format(
        title=title,
        chart_title=title,
        width=w,
        height=h,
        inner_height=h - 60,
        option_json=json.dumps(option, ensure_ascii=False)
    )


# ─────────────────────────────────────────────
# Phase 2: ECharts option → pyecharts 代码 (LLM辅助)
# ─────────────────────────────────────────────

PYECHARTS_TEMPLATE = '''#!/usr/bin/env python3
"""
Auto-generated pyecharts chart: {title}
Source: {source_file}
"""
import json, pathlib
from pyecharts.charts import {chart_class}
from pyecharts import options as opts
from pyecharts.globals import ThemeType

BASE = pathlib.Path(__file__).resolve().parent
OUT_DIR = BASE / "dist"
OUT_DIR.mkdir(parents=True, exist_ok=True)

def create_chart():
    """生成图表，返回 Chart 对象"""
{code_body}
    return chart

if __name__ == "__main__":
    chart = create_chart()
    # 渲染为HTML（可浏览器查看）
    html_path = OUT_DIR / "{safe_title}_pyecharts.html"
    chart.render(str(html_path))
    print(f"pyecharts HTML: {{html_path}}")
    
    # 也输出option JSON供比对
    opt_path = OUT_DIR / "{safe_title}_option.json"
    opt_path.write_text(
        json.dumps(chart.get_options(), ensure_ascii=False, indent=2),
        encoding="utf-8"
    )
    print(f"option JSON: {{opt_path}}")
'''

# ECharts type → pyecharts class mapping
CHART_TYPE_MAP = {
    "line": "Line",
    "bar": "Bar",
    "pie": "Pie",
    "scatter": "Scatter",
    "radar": "Radar",
    "funnel": "Funnel",
    "gauge": "Gauge",
    "heatmap": "HeatMap",
    "graph": "Graph",
    "tree": "Tree",
    "treemap": "Treemap",
    "sankey": "Sankey",
    "boxplot": "Boxplot",
    "candlestick": "Candlestick",
    "effectScatter": "EffectScatter",
    "lines": "Lines",
    "map": "Map",
    "parallel": "Parallel",
    "picture": "Picture",
    "pictorialBar": "PictorialBar",
    "sunburst": "Sunburst",
    "themeRiver": "ThemeRiver",
    "wordCloud": "WordCloud",
    "custom": "Custom",
}


def _infer_chart_type(option):
    """从ECharts option推断图表类型"""
    series_list = option.get("series", [])
    if not series_list:
        return "Bar"
    st = series_list[0].get("type", "bar")
    return CHART_TYPE_MAP.get(st, "Bar")


def _build_pyecharts_code_simple(option, chart_type):
    """
    简单转换：支持常见图表类型的ECharts option → pyecharts代码。
    仅处理简单情况，复杂图表走LLM。
    """
    series = option.get("series", [])
    if not series:
        return "# 无series数据"
    
    lines = []
    indent = "    "
    chart_class = CHART_TYPE_MAP.get(chart_type, "Bar")
    
    # 获取xAxis数据 (category)
    x_data = None
    xaxis = option.get("xAxis", {})
    if isinstance(xaxis, dict) and xaxis.get("type") == "category":
        x_data = xaxis.get("data", [])
    elif isinstance(xaxis, list) and len(xaxis) > 0:
        x_data = xaxis[0].get("data", []) if isinstance(xaxis[0], dict) else None
    
    # 构建pyecharts链式调用
    if chart_class == "Pie":
        # 饼图特殊处理
        lines.append(f"chart = {chart_class}(init_opts=opts.InitOpts(theme=ThemeType.DARK))")
        lines.append(f"chart.add_dataset([], {{}})")  # placeholder
        
        # 处理饼图数据
        pie_data = series[0].get("data", [])
        data_pairs = []
        for item in pie_data:
            if isinstance(item, dict):
                data_pairs.append(f'{{"value": {item.get("value", 0)}, "name": "{item.get("name", "")}"}}')
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                data_pairs.append(f'{{"value": {item[1]}, "name": "{item[0]}"}}')
        
        if data_pairs:
            lines.append(f'chart.add_dataset([{", ".join(data_pairs)}], {{"source": "data"}})')
        
        lines.append(f"""chart.set_global_opts(
    title_opts=opts.TitleOpts(title="{option.get('title', {}).get('text', 'Chart')}"),
    legend_opts=opts.LegendOpts(pos_left="center"),
)""")
    elif chart_class in ("Line", "Bar", "Scatter"):
        lines.append(f"chart = {chart_class}(init_opts=opts.InitOpts(theme=ThemeType.DARK))")
        
        if x_data:
            x_str = json.dumps(x_data, ensure_ascii=False)
            lines.append(f"chart.add_xaxis({x_str})")
        
        for s in series:
            sname = s.get("name", "series")
            sdata = s.get("data", [])
            sdata_str = json.dumps(sdata, ensure_ascii=False)
            
            if chart_class == "Line":
                lines.append(f'chart.add_yaxis("{sname}", {sdata_str}, is_smooth=True, linestyle_opts=opts.LineStyleOpts(width=2))')
            elif chart_class == "Bar":
                lines.append(f'chart.add_yaxis("{sname}", {sdata_str})')
            else:
                lines.append(f'chart.add_yaxis("{sname}", {sdata_str})')
        
        # 全局选项
        title_text = ""
        if isinstance(option.get("title"), dict):
            title_text = option["title"].get("text", "")
        
        lines.append(f"""chart.set_global_opts(
    title_opts=opts.TitleOpts(title="{title_text}"),
    tooltip_opts=opts.TooltipOpts(trigger="axis"),
    legend_opts=opts.LegendOpts(pos_left="center"),
)""")
    else:
        # 其他类型用通用方法
        lines.append(f"chart = {chart_class}(init_opts=opts.InitOpts(theme=ThemeType.DARK))")
        lines.append("# TODO: 需要手动调整此处的数据映射")
        lines.append(f"chart.add_dataset({json.dumps(series, ensure_ascii=False)}, {{'source': 'data'}})")
    
    return "\n".join(lines)


def cmd_convert(args):
    """Phase 2: ECharts option JSON → pyecharts Python代码"""
    data = _load_json(args.input)
    out_dir = pathlib.Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    
    source_name = pathlib.Path(args.input).name
    
    if "slides" in data:
        outputs = []
        for i, slide in enumerate(data.get("slides", [])):
            if slide.get("type") == "chart" and "option" in slide:
                out = _generate_pyecharts_file(slide, source_name, out_dir, i)
                outputs.append(out)
        print(f"[OK] 生成 {len(outputs)} 个pyecharts脚本:")
        for p in outputs:
            print(f"   {p}")
        return outputs
    else:
        out = _generate_pyecharts_file(data, source_name, out_dir)
        print(f"[OK] pyecharts脚本: {out}")
        # 也生成HTML预览
        preview_files = cmd_preview(args)
        return [out]


def _generate_pyecharts_file(data, source_name, out_dir, idx=None):
    """生成单个pyecharts Python脚本"""
    option = data.get("option", {})
    if not option:
        raise ValueError("JSON中缺少 'option' 字段 (ECharts option)")
    
    chart_type = _infer_chart_type(option)
    chart_class = CHART_TYPE_MAP.get(chart_type, "Bar")
    title = data.get("chart_title", data.get("slide_title", f"chart_{idx}" if idx is not None else "chart"))
    safe_title = _safe_name(title)
    
    # 生成pyecharts代码
    code_body = _build_pyecharts_code_simple(option, chart_type)
    
    code = PYECHARTS_TEMPLATE.format(
        title=title,
        source_file=source_name,
        chart_class=chart_class,
        code_body=code_body,
        safe_title=safe_title
    )
    
    # 写入文件
    suffix = f"_{idx}" if idx is not None else ""
    out_path = out_dir / f"{safe_title}{suffix}_pyecharts.py"
    out_path.write_text(code, encoding="utf-8")
    
    return out_path


# ─────────────────────────────────────────────
# Phase 3: pyecharts → PNG截图 → PPT组装
# ─────────────────────────────────────────────

def cmd_build(args):
    """Phase 3: 生成PPT"""
    data = _load_json(args.input)
    out_dir = pathlib.Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    
    if "slides" not in data:
        # 单图表模式：先转成presentation格式
        data = {
            "title": data.get("chart_title", data.get("slide_title", "Presentation")),
            "theme": {},
            "slides": [data]
        }
    
    # Step 3a: 为所有chart slide生成HTML并截图
    chart_images = _render_charts_to_images(data, out_dir)
    
    # Step 3b: 组装PPT
    ppt_path = _build_pptx(data, chart_images, out_dir)
    
    print(f"[OK] PPT生成: {ppt_path}")
    return ppt_path


def _render_charts_to_images(data, out_dir):
    """使用playwright将ECharts HTML渲染为PNG截图"""
    from playwright.sync_api import sync_playwright
    
    images = {}  # slide_index → image_path
    temp_dir = out_dir / "_temp_html"
    temp_dir.mkdir(parents=True, exist_ok=True)
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        
        for i, slide in enumerate(data.get("slides", [])):
            if slide.get("type") != "chart" or "option" not in slide:
                continue
            
            # 生成临时HTML
            html = _render_chart_html(slide, data.get("theme", {}))
            temp_html = temp_dir / f"_chart_{i}.html"
            temp_html.write_text(html, encoding="utf-8")
            
            # 截图
            try:
                page.goto(temp_html.resolve().as_uri(), wait_until="networkidle", timeout=15000)
                time.sleep(0.5)  # 等待ECharts渲染
                
                img_path = out_dir / f"chart_{i}.png"
                page.screenshot(path=str(img_path), full_page=False)
                images[i] = img_path
                print(f"  📸 chart_{i}.png ({img_path.stat().st_size / 1024:.1f} KB)")
            except Exception as e:
                print(f"  [WARN] chart_{i} 截图失败: {e}")
        
        browser.close()
    
    # 清理临时HTML
    shutil.rmtree(temp_dir, ignore_errors=True)
    
    return images


def _build_pptx(data, chart_images, out_dir):
    """组装PPTX文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    
    theme = data.get("theme", {})
    primary = _hex_to_rgb(theme.get("primary", "#4369e8"))
    secondary = _hex_to_rgb(theme.get("secondary", "#2ecca8"))
    accent = _hex_to_rgb(theme.get("accent", "#f39c12"))
    dark = RGBColor(0x1A, 0x1A, 0x2E)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    gray = RGBColor(0x88, 0x99, 0xBB)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def _add_text(slide, x, y, w, h, text, size=14, color=white, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Emu(int(x * 914400)), Emu(int(y * 914400)),
            Emu(int(w * 914400)), Emu(int(h * 914400))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return tf
    
    def _add_rect(slide, x, y, w, h, color):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x * 914400)), Emu(int(y * 914400)),
            Emu(int(w * 914400)), Emu(int(h * 914400))
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        return sh
    
    for i, slide_data in enumerate(data.get("slides", [])):
        sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _add_bg(sl, dark)
        
        stype = slide_data.get("type", "content")
        
        if stype == "cover":
            # 封面
            title = slide_data.get("title", "")
            subtitle = slide_data.get("subtitle", "")
            _add_text(sl, 1.5, 1.8, 10, 1.2, title, 40, white, True, PP_ALIGN.CENTER)
            if subtitle:
                _add_text(sl, 1.5, 3.0, 10, 0.6, subtitle, 16, gray, False, PP_ALIGN.CENTER)
            
            # 指标
            metrics = slide_data.get("metrics", [])
            for j, m in enumerate(metrics):
                x = 2.0 + j * 3.5
                _add_rect(sl, x, 4.0, 2.8, 1.5, RGBColor(0x1F, 0x1F, 0x35))
                _add_text(sl, x + 0.2, 4.2, 2.4, 0.5, m.get("value", ""), 28, primary, True, PP_ALIGN.CENTER)
                _add_text(sl, x + 0.2, 4.8, 2.4, 0.4, m.get("label", ""), 12, gray, False, PP_ALIGN.CENTER)
        
        elif stype == "chart":
            # 图表幻灯片
            title = slide_data.get("title", slide_data.get("chart_title", "Chart"))
            _add_text(sl, 1.0, 0.4, 11, 0.6, title, 24, white, True)
            
            # 插入图表截图
            if i in chart_images:
                img_path = chart_images[i]
                sl.shapes.add_picture(
                    str(img_path),
                    Emu(int(1.5 * 914400)), Emu(int(1.2 * 914400)),
                    Emu(int(10 * 914400)), Emu(int(5.5 * 914400))
                )
            else:
                _add_text(sl, 3.0, 3.0, 7, 1, "[图表渲染失败]", 18, gray, False, PP_ALIGN.CENTER)
            
            # 描述
            desc = slide_data.get("description", "")
            if desc:
                _add_text(sl, 1.0, 6.8, 11, 0.4, desc, 11, gray)
        
        elif stype == "content":
            # 内容幻灯片
            title = slide_data.get("title", "")
            key_point = slide_data.get("key_point", "")
            _add_text(sl, 1.0, 0.4, 11, 0.6, title, 24, white, True)
            if key_point:
                _add_text(sl, 1.0, 1.0, 11, 0.5, key_point, 13, gray)
            
            cards = slide_data.get("cards", [])
            for j, card in enumerate(cards):
                cy = 1.8 + j * 1.8
                _add_rect(sl, 1.0, cy, 11.3, 1.5, RGBColor(0x1F, 0x1F, 0x35))
                _add_text(sl, 1.3, cy + 0.1, 10, 0.3, card.get("title", ""), 14, white, True)
                for k, item in enumerate(card.get("items", [])[:5]):
                    _add_text(sl, 1.5, cy + 0.5 + k * 0.2, 10, 0.2, f"  {item}", 10, secondary)
        
        elif stype == "end":
            # 结束页
            title = slide_data.get("title", "Thank You")
            _add_text(sl, 1.0, 3.0, 11, 1, title, 36, white, True, PP_ALIGN.CENTER)
            sub = slide_data.get("subtitle", "")
            if sub:
                _add_text(sl, 1.0, 4.0, 11, 0.5, sub, 14, gray, False, PP_ALIGN.CENTER)
        
        else:
            # 默认内容页
            title = slide_data.get("title", "")
            _add_text(sl, 1.0, 0.4, 11, 0.6, title, 24, white, True)
    
    ppt_path = out_dir / f"{_safe_name(data.get('title', 'presentation'))}.pptx"
    prs.save(str(ppt_path))
    return ppt_path


# ─────────────────────────────────────────────
# 完整管道: run
# ─────────────────────────────────────────────

def cmd_run(args):
    """完整三阶段管道 + HITL"""
    data = _load_json(args.input)
    dist = pathlib.Path(args.out_dir)
    dist.mkdir(parents=True, exist_ok=True)
    
    print("=" * 60)
    print("  ECharts → pyecharts → PPT 管道")
    print(f"  输入: {args.input}")
    print("=" * 60)
    
    # ── Phase 1: HTML预览 ──
    print("\n📌 Phase 1: 生成预览HTML")
    print("-" * 40)
    preview_files = cmd_preview(args)
    
    print(f"\n[PAUSE]  [Gate 1] 请在浏览器中打开以上HTML文件，确认图表效果")
    print(f"   确认后输入 y 继续，输入 n 中止，或输入修改建议:")
    
    if not args.yes:
        resp = input(">> ").strip().lower()
        if resp.startswith("n"):
            print("⛔ 管道中止 (Gate 1)")
            return
        elif resp and resp != "y":
            print(f"📝 记录修改建议: {resp}")
            print("[PAUSE]  请修改JSON后重新运行，或输入 y 继续使用当前版本...")
            resp2 = input(">> ").strip().lower()
            if resp2.startswith("n"):
                print("⛔ 管道中止")
                return
    
    # ── Phase 2: 代码转换 ──
    print("\n📌 Phase 2: 转换为pyecharts代码")
    print("-" * 40)
    py_files = cmd_convert(args)
    
    print(f"\n[PAUSE]  [Gate 2] 请检查生成的pyecharts脚本:")
    for pf in py_files:
        print(f"   📄 {pf}")
    print(f"   确认后输入 y 继续构建PPT，或输入 n 中止:")
    
    if not args.yes:
        resp = input(">> ").strip().lower()
        if resp.startswith("n"):
            print("⛔ 管道中止 (Gate 2)")
            return
    
    # ── Phase 3: PPT构建 ──
    print("\n📌 Phase 3: 构建PPT")
    print("-" * 40)
    ppt_path = cmd_build(args)
    
    print(f"\n[OK] 管道完成!")
    print(f"   📊 预览HTML: {dist / '*.html'}")
    print(f"   🐍 pyecharts: {dist / '*_pyecharts.py'}")
    print(f"   📑 PPT: {ppt_path}")


# ─────────────────────────────────────────────
# 工具函数
# ─────────────────────────────────────────────

def _load_json(path):
    """加载JSON文件"""
    p = pathlib.Path(path)
    if not p.exists():
        # 尝试相对路径
        p = BASE / path
    if not p.exists():
        print(f"[FAIL] 文件不存在: {path}")
        sys.exit(1)
    return json.loads(p.read_text(encoding="utf-8"))


def _safe_name(name):
    """文件名安全化"""
    safe = re.sub(r'[^\w\u4e00-\u9fff\-]', '_', str(name))
    return safe.strip('_') or 'chart'


def _hex_to_rgb(hex_str):
    """#RRGGBB -> RGBColor"""
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip('#')
    if len(h) < 6:
        h = h * 6
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─────────────────────────────────────────────
# CLI入口
# ─────────────────────────────────────────────

def _add_common(sub):
    """Add common arguments to a subparser"""
    sub.add_argument("--out-dir", default=str(DIST_DIR),
                     help="输出目录 (默认: ./dist)")
    sub.add_argument("--yes", "-y", action="store_true",
                     help="跳过所有Gate确认 (全自动模式)")
    return sub

def main():
    parser = argparse.ArgumentParser(
        description="ECharts → pyecharts → PPT 三阶段管道",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    
    sub = parser.add_subparsers(dest="command", required=True)
    
    # preview
    p_preview = _add_common(sub.add_parser("preview", help="生成ECharts HTML预览"))
    p_preview.add_argument("input", help="ECharts option JSON文件")
    
    # convert
    p_convert = _add_common(sub.add_parser("convert", help="转换为pyecharts代码"))
    p_convert.add_argument("input", help="ECharts option JSON文件")
    
    # build
    p_build = _add_common(sub.add_parser("build", help="构建PPT"))
    p_build.add_argument("input", help="Presentation JSON文件")
    
    # run
    p_run = _add_common(sub.add_parser("run", help="完整管道 (preview→convert→build + Gate)"))
    p_run.add_argument("input", help="Presentation JSON文件")
    
    args = parser.parse_args()
    
    # 分发
    if args.command == "preview":
        cmd_preview(args)
    elif args.command == "convert":
        cmd_convert(args)
    elif args.command == "build":
        cmd_build(args)
    elif args.command == "run":
        cmd_run(args)


if __name__ == "__main__":
    main()
